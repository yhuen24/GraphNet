"""
Document processing module for GraphNet.
Handles extraction of text from various document formats.
"""

import os
from typing import List, Dict, Any
from io import BytesIO
import logging
import json

# Document processing libraries
import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation
import pandas as pd

from config import config

# Setup logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


# =============================================================================
# Chunking defaults
# -----------------------------------------------------------------------------
# These are used when config.CHUNK_SIZE / config.CHUNK_OVERLAP are missing OR
# when the configured values are clearly too small for modern LLMs (Gemini,
# Claude, GPT-4o all handle 20k+ char chunks comfortably).
# =============================================================================

DEFAULT_CHUNK_SIZE = 6000        # ~1.5k tokens - forces dense extraction
DEFAULT_CHUNK_OVERLAP = 1000     # Gives enough overlap to catch boundary relationships
MIN_REASONABLE_CHUNK_SIZE = 2000


class DocumentProcessor:
    """Process various document formats and extract text"""

    @staticmethod
    def process_file(file_path: str = None, file_bytes: bytes = None,
                     file_extension: str = None, filename: str = "unknown") -> Dict[str, Any]:
        """
        Process a file and extract text content

        Args:
            file_path: Path to the file (for local files)
            file_bytes: Bytes content of the file (for uploaded files)
            file_extension: Extension of the file
            filename: Name of the file

        Returns:
            Dictionary containing extracted text and metadata
        """
        try:

            if file_path:
                file_extension = os.path.splitext(file_path)[1].lower()

            if not file_extension:
                raise ValueError("File extension not provided")

            if file_extension in ['.txt', '.md']:
                return DocumentProcessor._process_text(file_path, file_bytes, filename)
            elif file_extension == '.pdf':
                return DocumentProcessor._process_pdf(file_path, file_bytes, filename)
            elif file_extension == '.docx':
                return DocumentProcessor._process_docx(file_path, file_bytes, filename)
            elif file_extension == '.xlsx':
                return DocumentProcessor._process_xlsx(file_path, file_bytes, filename)
            elif file_extension == '.pptx':
                return DocumentProcessor._process_pptx(file_path, file_bytes, filename)
            elif file_extension == '.csv':
                return DocumentProcessor._process_csv(file_path, file_bytes, filename)
            elif file_extension == '.json':
                return DocumentProcessor._process_json(file_path, file_bytes, filename)
            else:
                raise ValueError(f"Unsupported file format: {file_extension}")

        except Exception as e:
            logger.error(f"Error processing file: {str(e)}")
            return {
                "success": False,
                "error": str(e),
                "text": "",
                "metadata": {}
            }

    @staticmethod
    def _process_text(file_path: str = None, file_bytes: bytes = None, filename: str = "unknown") -> Dict[str, Any]:
        try:
            if file_bytes:
                text = file_bytes.decode('utf-8')
            else:
                with open(file_path, 'r', encoding='utf-8') as f:
                    text = f.read()

            return {
                "success": True,
                "text": text,
                "metadata": {
                    "format": "text",
                    "filename": filename,
                    "character_count": len(text),
                    "word_count": len(text.split())
                }
            }
        except Exception as e:
            logger.error(f"Error processing text file: {str(e)}")
            raise

    @staticmethod
    def _process_pdf(file_path: str = None, file_bytes: bytes = None, filename: str = "unknown") -> Dict[str, Any]:
        try:
            text = ""
            source = BytesIO(file_bytes) if file_bytes else file_path

            with pdfplumber.open(source) as pdf:
                page_count = len(pdf.pages)
                for i, page in enumerate(pdf.pages):
                    try:
                        page_text = page.extract_text()
                        if page_text:
                            text += f"\n--- Page {i + 1} ---\n"
                            text += page_text
                    except Exception as e:
                        logger.warning(f"Could not extract text from page {i + 1}: {str(e)}")
                        continue

            if not text.strip():
                raise ValueError("No text could be extracted from PDF")

            return {
                "success": True,
                "text": text,
                "metadata": {
                    "format": "pdf",
                    "filename": filename,
                    "page_count": page_count,
                    "character_count": len(text),
                    "word_count": len(text.split())
                }
            }
        except Exception as e:
            logger.error(f"Error processing PDF: {str(e)}")
            raise

    @staticmethod
    def _process_docx(file_path: str = None, file_bytes: bytes = None, filename: str = "unknown") -> Dict[str, Any]:
        try:
            if file_bytes:
                doc = DocxDocument(BytesIO(file_bytes))
            else:
                doc = DocxDocument(file_path)

            text = ""
            for para in doc.paragraphs:
                text += para.text + "\n"

            # Extract tables
            for table in doc.tables:
                text += "\n--- Table ---\n"
                for row in table.rows:
                    row_text = " | ".join([cell.text for cell in row.cells])
                    text += row_text + "\n"

            return {
                "success": True,
                "text": text,
                "metadata": {
                    "format": "docx",
                    "filename": filename,
                    "paragraph_count": len(doc.paragraphs),
                    "table_count": len(doc.tables),
                    "character_count": len(text),
                    "word_count": len(text.split())
                }
            }
        except Exception as e:
            logger.error(f"Error processing DOCX: {str(e)}")
            raise

    @staticmethod
    def _process_xlsx(file_path: str = None, file_bytes: bytes = None, filename: str = "unknown") -> Dict[str, Any]:
        try:
            if file_bytes:
                df_dict = pd.read_excel(BytesIO(file_bytes), sheet_name=None)
            else:
                df_dict = pd.read_excel(file_path, sheet_name=None)

            text = ""
            for sheet_name, df in df_dict.items():
                text += f"\n--- Sheet: {sheet_name} ---\n"
                text += df.to_string(index=False)
                text += "\n\n"

            return {
                "success": True,
                "text": text,
                "metadata": {
                    "format": "xlsx",
                    "filename": filename,
                    "sheet_count": len(df_dict),
                    "character_count": len(text),
                    "word_count": len(text.split())
                }
            }
        except Exception as e:
            logger.error(f"Error processing XLSX: {str(e)}")
            raise

    @staticmethod
    def _process_pptx(file_path: str = None, file_bytes: bytes = None, filename: str = "unknown") -> Dict[str, Any]:
        try:
            if file_bytes:
                prs = Presentation(BytesIO(file_bytes))
            else:
                prs = Presentation(file_path)

            text = ""
            for slide_num, slide in enumerate(prs.slides):
                text += f"\n--- Slide {slide_num + 1} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"

            return {
                "success": True,
                "text": text,
                "metadata": {
                    "format": "pptx",
                    "filename": filename,
                    "slide_count": len(prs.slides),
                    "character_count": len(text),
                    "word_count": len(text.split())
                }
            }
        except Exception as e:
            logger.error(f"Error processing PPTX: {str(e)}")
            raise

    @staticmethod
    def _process_csv(file_path: str = None, file_bytes: bytes = None, filename: str = "unknown") -> Dict[str, Any]:
        try:
            if file_bytes:
                df = pd.read_csv(BytesIO(file_bytes))
            else:
                df = pd.read_csv(file_path)

            text = df.to_string(index=False)

            return {
                "success": True,
                "text": text,
                "metadata": {
                    "format": "csv",
                    "filename": filename,
                    "row_count": len(df),
                    "column_count": len(df.columns),
                    "character_count": len(text),
                    "word_count": len(text.split())
                }
            }
        except Exception as e:
            logger.error(f"Error processing CSV: {str(e)}")
            raise

    @staticmethod
    def _process_json(file_path: str = None, file_bytes: bytes = None, filename: str = "unknown") -> Dict[str, Any]:
        try:
            if file_bytes:
                data = json.loads(file_bytes.decode('utf-8'))
            else:
                with open(file_path, 'r', encoding='utf-8') as f:
                    data = json.load(f)

            text = json.dumps(data, indent=2)

            return {
                "success": True,
                "text": text,
                "metadata": {
                    "format": "json",
                    "filename": filename,
                    "character_count": len(text),
                    "word_count": len(text.split())
                }
            }
        except Exception as e:
            logger.error(f"Error processing JSON: {str(e)}")
            raise

    # -------------------------------------------------------------------------
    # Chunking
    # -------------------------------------------------------------------------

    @staticmethod
    def chunk_text(text: str, chunk_size: int = None, chunk_overlap: int = None) -> List[str]:
        """
        Split text into chunks for LLM processing.

        Uses sentence-boundary-aware splitting so entities and relationships
        aren't cut mid-sentence. Falls back gracefully when no sentence break
        is available near the target boundary.

        Defaults are tuned for modern long-context LLMs (Gemini, Claude, GPT-4o).
        If config.CHUNK_SIZE is set but is too small to be efficient
        (< 8,000 chars), it is overridden with the modern default and a warning
        is logged — otherwise a 14-page PDF ends up as 27+ micro-chunks and
        extraction takes hours.

        Args:
            text: Text to chunk
            chunk_size: Maximum characters per chunk (default: 40,000)
            chunk_overlap: Characters of overlap between consecutive chunks
                           (default: 2,000)

        Returns:
            List of text chunks
        """
        # ---- Resolve chunk size ----
        if chunk_size is None:
            configured = getattr(config, "CHUNK_SIZE", None)
            if configured and configured >= MIN_REASONABLE_CHUNK_SIZE:
                chunk_size = configured
            else:
                if configured and configured < MIN_REASONABLE_CHUNK_SIZE:
                    logger.warning(
                        f"config.CHUNK_SIZE={configured} is too small for efficient "
                        f"LLM extraction. Using {DEFAULT_CHUNK_SIZE} instead. "
                        f"Update config.CHUNK_SIZE to silence this warning."
                    )
                chunk_size = DEFAULT_CHUNK_SIZE

        # ---- Resolve overlap ----
        if chunk_overlap is None:
            configured_overlap = getattr(config, "CHUNK_OVERLAP", None)
            # Overlap should be ~5% of chunk size and never exceed it
            if configured_overlap and 0 < configured_overlap < chunk_size:
                chunk_overlap = configured_overlap
            else:
                chunk_overlap = DEFAULT_CHUNK_OVERLAP

        # ---- Short-circuit: text fits in a single chunk ----
        if len(text) <= chunk_size:
            logger.info(f"Text ({len(text)} chars) fits in a single chunk")
            return [text]

        # ---- Sentence-boundary-aware chunking ----
        chunks = []
        start = 0
        text_length = len(text)

        while start < text_length:
            end = min(start + chunk_size, text_length)

            # Try to break at a sentence boundary within the last 20% of the chunk
            if end < text_length:
                search_start = max(start, end - int(chunk_size * 0.2))
                last_period = text.rfind('. ', search_start, end)
                last_newline = text.rfind('\n', search_start, end)
                break_point = max(last_period, last_newline)
                if break_point > search_start:
                    end = break_point + 1

            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)

            # Advance, with overlap
            next_start = end - chunk_overlap
            # Guard against infinite loop if overlap >= progress
            if next_start <= start:
                next_start = end
            start = next_start

            if start >= text_length:
                break

        logger.info(
            f"Split text ({text_length} chars) into {len(chunks)} chunks "
            f"(chunk_size={chunk_size}, overlap={chunk_overlap})"
        )
        return chunks